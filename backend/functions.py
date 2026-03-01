# functions.py - User-Defined Functions (UDFs) for the Pixeltable Agent
import os
import traceback
from datetime import datetime
from typing import Optional, Dict, Any, List, Union

import requests
import yfinance as yf
from duckduckgo_search import DDGS
import pixeltable as pxt


@pxt.udf
def get_latest_news(topic: str) -> str:
    """Fetch latest news for a given topic using NewsAPI."""
    try:
        api_key = os.environ.get("NEWS_API_KEY")
        if not api_key:
            return "Error: NewsAPI key not found in environment variables."

        url = "https://newsapi.org/v2/everything"
        params = {
            "q": topic,
            "apiKey": api_key,
            "sortBy": "publishedAt",
            "language": "en",
            "pageSize": 5,
        }

        response = requests.get(url, params=params, timeout=10)

        if response.status_code != 200:
            return f"Error: NewsAPI request failed ({response.status_code}): {response.text}"

        data = response.json()
        articles = data.get("articles", [])

        if not articles:
            return f"No recent news found for '{topic}'."

        formatted_news = []
        for i, article in enumerate(articles[:3], 1):
            pub_date = datetime.fromisoformat(
                article["publishedAt"].replace("Z", "+00:00")
            ).strftime("%Y-%m-%d")
            formatted_news.append(
                f"{i}. [{pub_date}] {article['title']}\n   {article['description']}"
            )

        return "\n\n".join(formatted_news)

    except requests.Timeout:
        return "Error: NewsAPI request timed out."
    except requests.RequestException as e:
        return f"Error making NewsAPI request: {str(e)}"
    except Exception as e:
        return f"Unexpected error fetching news: {str(e)}."


@pxt.udf
def search_news(keywords: str, max_results: int = 5) -> str:
    """Search news using DuckDuckGo and return results."""
    try:
        with DDGS() as ddgs:
            results = list(
                ddgs.news(
                    keywords=keywords,
                    region="wt-wt",
                    safesearch="off",
                    timelimit="m",
                    max_results=max_results,
                )
            )
            if not results:
                return "No news results found."

            formatted_results = []
            for i, r in enumerate(results, 1):
                formatted_results.append(
                    f"{i}. Title: {r.get('title', 'N/A')}\n"
                    f"   Source: {r.get('source', 'N/A')}\n"
                    f"   Published: {r.get('date', 'N/A')}\n"
                    f"   URL: {r.get('url', 'N/A')}\n"
                    f"   Snippet: {r.get('body', 'N/A')}\n"
                )
            return "\n".join(formatted_results)
    except Exception as e:
        print(f"DuckDuckGo search failed: {str(e)}")
        return f"Search failed: {str(e)}."


@pxt.udf
def fetch_financial_data(ticker: str) -> str:
    """Fetch financial summary data for a given company ticker using yfinance."""
    try:
        if not ticker:
            return "Error: No ticker symbol provided."

        stock = yf.Ticker(ticker)
        info = stock.info

        if not info or info.get("quoteType") == "MUTUALFUND":
            hist = stock.history(period="1d")
            if hist.empty:
                return f"Error: No data found for ticker '{ticker}'. It might be delisted or incorrect."
            else:
                return f"Limited info for '{ticker}'. Previous Close: {hist['Close'].iloc[-1]:.2f} (if available)."

        data_points = {
            "Company Name": info.get("shortName") or info.get("longName"),
            "Symbol": info.get("symbol"),
            "Exchange": info.get("exchange"),
            "Quote Type": info.get("quoteType"),
            "Currency": info.get("currency"),
            "Current Price": info.get("currentPrice") or info.get("regularMarketPrice") or info.get("bid"),
            "Previous Close": info.get("previousClose"),
            "Open": info.get("open"),
            "Day Low": info.get("dayLow"),
            "Day High": info.get("dayHigh"),
            "Volume": info.get("volume") or info.get("regularMarketVolume"),
            "Market Cap": info.get("marketCap"),
            "Trailing P/E": info.get("trailingPE"),
            "Forward P/E": info.get("forwardPE"),
            "Dividend Yield": info.get("dividendYield"),
            "52 Week Low": info.get("fiftyTwoWeekLow"),
            "52 Week High": info.get("fiftyTwoWeekHigh"),
            "Avg Volume (10 day)": info.get("averageDailyVolume10Day"),
        }

        formatted_data = [
            f"Financial Summary for {data_points.get('Company Name', ticker)} ({data_points.get('Symbol', ticker).upper()}) - {data_points.get('Quote Type', 'N/A')}"
        ]
        formatted_data.append("-" * 40)

        for key, value in data_points.items():
            if value is not None:
                formatted_value = value
                if key in ["Current Price", "Previous Close", "Open", "Day Low", "Day High", "52 Week Low", "52 Week High"] and isinstance(value, (int, float)):
                    formatted_value = f"{value:.2f} {data_points.get('Currency', '')}".strip()
                elif key in ["Volume", "Market Cap", "Avg Volume (10 day)"] and isinstance(value, (int, float)):
                    if value > 1_000_000_000:
                        formatted_value = f"{value / 1_000_000_000:.2f}B"
                    elif value > 1_000_000:
                        formatted_value = f"{value / 1_000_000:.2f}M"
                    elif value > 1_000:
                        formatted_value = f"{value / 1_000:.2f}K"
                    else:
                        formatted_value = f"{value:,}"
                elif key == "Dividend Yield" and isinstance(value, (int, float)):
                    formatted_value = f"{value * 100:.2f}%"
                elif key in ("Trailing P/E", "Forward P/E") and isinstance(value, (int, float)):
                    formatted_value = f"{value:.2f}"

                formatted_data.append(f"{key}: {formatted_value}")

        try:
            latest_financials = stock.financials.iloc[:, 0]
            revenue = latest_financials.get("Total Revenue")
            net_income = latest_financials.get("Net Income")
            if revenue is not None or net_income is not None:
                formatted_data.append("-" * 40)
                fin_date = latest_financials.name.strftime("%Y-%m-%d")
                if revenue:
                    formatted_data.append(f"Latest Revenue ({fin_date}): ${revenue / 1e6:.2f}M")
                if net_income:
                    formatted_data.append(f"Latest Net Income ({fin_date}): ${net_income / 1e6:.2f}M")
        except Exception:
            pass

        return "\n".join(formatted_data)

    except Exception as e:
        traceback.print_exc()
        return f"Error fetching financial data for {ticker}: {str(e)}."


# ── Notification UDFs ─────────────────────────────────────────────────────────
# Each wraps a simple HTTP POST. Registered as agent tools so the chat agent
# can send notifications on the user's behalf ("summarize my docs and post to Slack").


@pxt.udf
def send_slack_message(message: str) -> str:
    """Send a message to a configured Slack channel via incoming webhook."""
    webhook_url = os.environ.get("SLACK_WEBHOOK_URL", "")
    if not webhook_url:
        return "Error: SLACK_WEBHOOK_URL not configured."
    try:
        resp = requests.post(webhook_url, json={"text": message}, timeout=10)
        if resp.status_code == 200:
            return f"Slack message sent successfully."
        return f"Slack error ({resp.status_code}): {resp.text}"
    except requests.RequestException as e:
        return f"Slack request failed: {e}"


@pxt.udf
def send_discord_message(message: str) -> str:
    """Send a message to a configured Discord channel via webhook."""
    webhook_url = os.environ.get("DISCORD_WEBHOOK_URL", "")
    if not webhook_url:
        return "Error: DISCORD_WEBHOOK_URL not configured."
    try:
        resp = requests.post(webhook_url, json={"content": message}, timeout=10)
        if resp.status_code in (200, 204):
            return f"Discord message sent successfully."
        return f"Discord error ({resp.status_code}): {resp.text}"
    except requests.RequestException as e:
        return f"Discord request failed: {e}"


@pxt.udf
def send_webhook(message: str, url: str = "") -> str:
    """POST a JSON payload to any webhook URL. Connects to n8n, Zapier, Make, or custom endpoints."""
    target_url = url or os.environ.get("WEBHOOK_URL", "")
    if not target_url:
        return "Error: No webhook URL provided and WEBHOOK_URL not configured."
    try:
        payload = {"text": message, "source": "pixelbot", "timestamp": datetime.utcnow().isoformat()}
        resp = requests.post(target_url, json=payload, timeout=10)
        if resp.status_code < 300:
            return f"Webhook delivered ({resp.status_code})."
        return f"Webhook error ({resp.status_code}): {resp.text}"
    except requests.RequestException as e:
        return f"Webhook request failed: {e}"


@pxt.udf
def extract_document_text(doc: pxt.Document) -> Optional[str]:
    """Extract full text from a document file, truncated for LLM summarization."""
    import pathlib

    try:
        path = pathlib.Path(doc)
        ext = path.suffix.lower()
        text = ""

        if ext == ".pdf":
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                pages = [page.extract_text() or "" for page in pdf.pages]
                text = "\n".join(pages)

        elif ext in (".docx", ".doc"):
            import docx as docx_lib

            doc_obj = docx_lib.Document(str(path))
            text = "\n".join(p.text for p in doc_obj.paragraphs if p.text.strip())

        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation

            prs = Presentation(str(path))
            slide_texts: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
            text = "\n".join(slide_texts)

        elif ext in (".xlsx", ".xls"):
            import openpyxl

            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            rows_text: List[str] = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows_text.append(", ".join(cells))
            wb.close()
            text = "\n".join(rows_text)

        elif ext in (".csv",):
            import csv

            with open(str(path), newline="", errors="ignore") as f:
                reader = csv.reader(f)
                csv_rows: List[str] = []
                for row in reader:
                    csv_rows.append(", ".join(row))
                    if len(csv_rows) >= 500:
                        break
            text = "\n".join(csv_rows)

        elif ext in (".txt", ".md", ".html", ".xml", ".rtf"):
            text = path.read_text(errors="ignore")

        else:
            text = path.read_text(errors="ignore")

        if not text or not text.strip():
            return None

        max_chars = 15_000
        if len(text) > max_chars:
            text = text[:max_chars] + "\n...[truncated]"

        return text.strip()

    except Exception as e:
        return f"[Error extracting text: {str(e)}]"


@pxt.udf
def build_tool_selection_messages(
    prompt: str,
    history_context: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    """Build the messages array for the tool-selection LLM call (Gemini format).
    Includes recent chat history so the model knows what 'that' or 'it' refers to."""
    msgs: List[Dict[str, Any]] = []
    if history_context:
        for item in reversed(history_context):
            role = item.get("role")
            content = item.get("content")
            if role and content:
                # Gemini uses "model" instead of "assistant"
                gemini_role = "model" if role == "assistant" else role
                msgs.append({"role": gemini_role, "parts": [{"text": content}]})
    msgs.append({"role": "user", "parts": [{"text": prompt}]})
    return msgs


@pxt.udf
def assemble_multimodal_context(
    question: str,
    tool_outputs: Optional[List[Dict[str, Any]]],
    doc_context: Optional[List[Union[Dict[str, Any], str]]],
    memory_context: Optional[List[Dict[str, Any]]] = None,
    chat_memory_context: Optional[List[Dict[str, Any]]] = None,
) -> str:
    """Construct a text block summarizing context types relevant to the user's question."""
    doc_context_str = "N/A"
    if doc_context:
        doc_items = []
        for item in doc_context:
            text = item.get("text", "") if isinstance(item, dict) else str(item)
            source = item.get("source_doc", "Unknown Document") if isinstance(item, dict) else "Unknown Document"
            source_name = os.path.basename(str(source))
            if text:
                doc_items.append(f"- [Source: {source_name}] {text}")
        if doc_items:
            doc_context_str = "\n".join(doc_items)

    memory_context_str = "N/A"
    if memory_context:
        memory_items = []
        for item in memory_context:
            content = item.get("content", "")
            item_type = item.get("type", "unknown")
            language = item.get("language")
            sim = item.get("sim")
            context_query = item.get("context_query", "Unknown Query")
            item_desc = f"""- [Memory Item | Type: {item_type}{f" ({language})" if language else ""} | Original Query: '{context_query}' {f"| Sim: {sim:.3f}" if sim is not None else ""}]
Content: {content[:100]}{"..." if len(content) > 100 else ""}"""
            memory_items.append(item_desc)
        if memory_items:
            memory_context_str = "\n".join(memory_items)

    chat_memory_context_str = "N/A"
    if chat_memory_context:
        chat_memory_items = []
        for item in chat_memory_context:
            content = item.get("content", "")
            role = item.get("role", "unknown")
            sim = item.get("sim")
            timestamp = item.get("timestamp")
            ts_str = timestamp.strftime("%Y-%m-%d %H:%M") if timestamp else "Unknown Time"
            item_desc = f"""- [Chat History | Role: {role} | Time: {ts_str} {f"| Sim: {sim:.3f}" if sim is not None else ""}]
Content: {content[:150]}{"..." if len(content) > 150 else ""}"""
            chat_memory_items.append(item_desc)
        if chat_memory_items:
            chat_memory_context_str = "\n".join(chat_memory_items)

    tool_outputs_str = str(tool_outputs) if tool_outputs else "N/A"

    text_content = f"""
ORIGINAL QUESTION:
{question}

AVAILABLE CONTEXT:

[TOOL RESULTS]
{tool_outputs_str}

[DOCUMENT CONTEXT]
{doc_context_str}

[MEMORY BANK CONTEXT]
{memory_context_str}

[CHAT HISTORY SEARCH CONTEXT] (Older messages relevant to the query)
{chat_memory_context_str}
"""
    return text_content.strip()


@pxt.udf
def assemble_final_messages(
    history_context: Optional[List[Dict[str, Any]]],
    multimodal_context_text: str,
    image_context: Optional[List[Dict[str, Any]]] = None,
    video_frame_context: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    """Construct the final list of messages for Gemini, incorporating all context types.

    Gemini format: role is "user"/"model", content is a list of parts
    ({"text": "..."} or {"inline_data": {"mime_type": ..., "data": base64}}).
    """
    messages = []

    if history_context:
        for item in reversed(history_context):
            role = item.get("role")
            content = item.get("content")
            if role and content:
                gemini_role = "model" if role == "assistant" else role
                messages.append({"role": gemini_role, "parts": [{"text": content}]})

    final_parts: List[Dict[str, Any]] = []

    if image_context:
        for item in image_context:
            if isinstance(item, dict) and "encoded_image" in item:
                image_data = item["encoded_image"]
                if isinstance(image_data, bytes):
                    image_data = image_data.decode("utf-8")
                elif not isinstance(image_data, str):
                    continue
                final_parts.append({
                    "inline_data": {"mime_type": "image/png", "data": image_data},
                })

    if video_frame_context:
        for item in video_frame_context:
            # search_video_frames returns "encoded_frame" key
            frame_data = None
            if isinstance(item, dict):
                frame_data = item.get("encoded_frame") or item.get("encoded_video_frame")
            if frame_data is None:
                continue
            if isinstance(frame_data, bytes):
                frame_data = frame_data.decode("utf-8")
            elif not isinstance(frame_data, str):
                continue
            final_parts.append({
                "inline_data": {"mime_type": "image/png", "data": frame_data},
            })

    final_parts.append({"text": multimodal_context_text})
    messages.append({"role": "user", "parts": final_parts})
    return messages


@pxt.udf
def assemble_follow_up_prompt(original_prompt: str, answer_text: str) -> str:
    """Construct the user content for the follow-up question LLM call.
    The system instruction lives in the Gemini config; this is just the user turn."""
    return f"Question: {original_prompt}\n\nAnswer: {answer_text}"
